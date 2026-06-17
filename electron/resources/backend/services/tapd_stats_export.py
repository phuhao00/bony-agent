"""TAPD 缺陷统计导出（MD / Excel / PDF / PPT）与 AI 整理报告。"""
from __future__ import annotations

import io
import json
import re
from datetime import datetime
from typing import Any

from utils.logger import setup_logger

logger = setup_logger("tapd_stats_export")

EXPORT_FORMATS = frozenset({"md", "excel", "pdf", "ppt"})

REPORT_CSS = """
body {
  font-family: -apple-system, BlinkMacSystemFont, "PingFang SC", "Microsoft YaHei", sans-serif;
  font-size: 11pt;
  line-height: 1.6;
  color: #1e293b;
}
h1 {
  font-size: 20pt;
  font-weight: 700;
  color: #0f172a;
  border-bottom: 2px solid #e2e8f0;
  padding-bottom: 10px;
  margin: 0 0 16px 0;
}
h2 {
  font-size: 14pt;
  font-weight: 600;
  color: #334155;
  margin: 20px 0 10px 0;
}
h3 {
  font-size: 12pt;
  font-weight: 600;
  color: #475569;
  margin: 14px 0 8px 0;
}
p { margin: 8px 0; }
ul, ol { margin: 8px 0 8px 22px; padding: 0; }
li { margin: 5px 0; }
strong { color: #0f172a; }
table {
  border-collapse: collapse;
  width: 100%;
  margin: 12px 0;
  font-size: 10pt;
}
th, td {
  border: 1px solid #e2e8f0;
  padding: 7px 10px;
  text-align: left;
}
th { background: #f8fafc; font-weight: 600; }
.meta {
  color: #64748b;
  font-size: 10pt;
  margin-bottom: 18px;
  line-height: 1.7;
}
.appendix {
  margin-top: 24px;
  padding-top: 12px;
  border-top: 1px solid #e2e8f0;
  color: #64748b;
  font-size: 9.5pt;
}
blockquote {
  border-left: 3px solid #6366f1;
  margin: 12px 0;
  padding: 4px 0 4px 14px;
  color: #475569;
}
hr { border: none; border-top: 1px solid #e2e8f0; margin: 18px 0; }
"""


def _range_label(stats: dict[str, Any]) -> str:
    rng = stats.get("range") or {}
    if rng.get("created_start") or rng.get("created_end"):
        return f"{rng.get('created_start') or '…'} ~ {rng.get('created_end') or '…'}"
    days = rng.get("days")
    if days in (None, 0):
        return "全部时间"
    return f"近 {days} 天"


def _compact_stats_for_ai(stats: dict[str, Any]) -> dict[str, Any]:
    """压缩统计数据供 AI 分析，避免塞入过长原始列表。"""
    return {
        "workspace_id": stats.get("workspace_id"),
        "range_label": _range_label(stats),
        "summary": stats.get("summary"),
        "by_status": stats.get("by_status"),
        "by_priority": stats.get("by_priority"),
        "by_owner_top10": (stats.get("by_owner") or [])[:10],
        "by_reporter_top10": (stats.get("by_reporter") or [])[:10],
        "recent_titles": [
            {
                "title": b.get("title"),
                "status": b.get("status_label") or b.get("status"),
                "priority": b.get("priority_label"),
                "owner": b.get("current_owner"),
            }
            for b in (stats.get("recent_bugs") or [])[:8]
        ],
    }


def _strip_ai_fences(text: str) -> str:
    t = (text or "").strip()
    if t.startswith("```"):
        t = re.sub(r"^```(?:markdown|md)?\s*", "", t)
        t = re.sub(r"\s*```$", "", t)
    return t.strip()


def _build_analysis_system_prompt(mode: str) -> str:
    if mode == "deep":
        return "\n".join(
            [
                "你是资深 QA 负责人，擅长把缺陷统计数据整理成可直接汇报的分析报告。",
                "请基于给定统计数据，输出一份中文 Markdown 报告正文（不要代码块包裹）。",
                "结构要求：",
                "1. ## 执行摘要 — 3~5 条要点，领导一眼能看懂",
                "2. ## 核心指标解读 — 用段落说明总量、开闭、未关闭占比，数字必须来自数据",
                "3. ## 分布与趋势 — 解读状态/优先级/人员分布，指出异常与集中度",
                "4. ## 风险与关注点 — 高优未关闭、积压、重复模块等",
                "5. ## 改进建议 — 3~5 条可执行建议，按优先级排序",
                "6. ## 跟进事项 — 如需跟进的清单",
                "禁止：粘贴 JSON、重复原始大表格、输出 Markdown 语法教学、编造数据。",
                "文风：专业、简洁、像写给团队的质量周报。",
            ]
        )
    return "\n".join(
        [
            "你是 QA 分析助手，请把 TAPD 缺陷统计数据整理成简明可读的中文 Markdown 报告正文。",
            "包含：执行摘要（要点列表）、核心发现（2~3 段）、主要风险、行动建议（2~4 条）。",
            "用叙述和要点呈现，不要大段原始表格，不要 JSON，不要编造。",
        ]
    )


def analyze_stats(
    stats: dict[str, Any],
    *,
    mode: str = "summary",
    user_note: str = "",
) -> dict[str, Any]:
    """AI 整理统计报告正文（Markdown）。"""
    analysis_mode = (mode or "summary").strip().lower()
    if analysis_mode not in {"summary", "deep"}:
        analysis_mode = "summary"

    compact = _compact_stats_for_ai(stats)
    try:
        stats_json = json.dumps(compact, ensure_ascii=False, indent=2)
    except TypeError:
        stats_json = str(compact)

    user_parts = [
        "以下是 TAPD 缺陷统计数据（JSON），请整理成分析报告正文：",
        stats_json,
    ]
    note = (user_note or "").strip()
    if note:
        user_parts.append(f"用户补充要求：{note}")

    system = _build_analysis_system_prompt(analysis_mode)
    user_message = "\n\n".join(user_parts)

    try:
        from langchain_core.messages import HumanMessage, SystemMessage

        from core.llm_provider import get_chat_llm

        llm = get_chat_llm(temperature=0.35)
        resp = llm.invoke(
            [SystemMessage(content=system), HumanMessage(content=user_message)]
        )
        analysis = _strip_ai_fences(str(getattr(resp, "content", "") or ""))
        if not analysis:
            return {"ok": False, "error": "模型未返回分析内容，请检查 LLM 配置"}
        logger.info("[tapd_export] AI report done mode=%s chars=%d", analysis_mode, len(analysis))
        return {"ok": True, "analysis": analysis, "mode": analysis_mode}
    except Exception as e:
        logger.error("[tapd_export] AI analysis failed: %s", e, exc_info=True)
        return {"ok": False, "error": f"AI 分析失败：{str(e)[:300]}"}


def _build_fallback_narrative(stats: dict[str, Any]) -> str:
    """无 AI 时的模板化可读报告（非原始表格 dump）。"""
    summary = stats.get("summary") or {}
    lines = [
        "## 执行摘要",
        "",
        f"- 统计范围 **{_range_label(stats)}**，范围内共 **{summary.get('total_in_range', '—')}** 条缺陷。",
        f"- 未关闭 **{summary.get('open', '—')}** 条，已关闭 **{summary.get('closed', '—')}** 条，未关闭占比 **{summary.get('open_rate', '—')}%**。",
        "",
        "## 分布概览",
        "",
    ]
    for title, key in [
        ("状态", "by_status"),
        ("优先级", "by_priority"),
        ("处理人 Top5", "by_owner"),
    ]:
        rows = stats.get(key) or []
        if not rows:
            continue
        lines.append(f"**{title}**：")
        for row in rows[:5]:
            lines.append(
                f"- {row.get('label')}：{row.get('count')} 条（{row.get('percent')}%）"
            )
        lines.append("")
    lines.extend(
        [
            "## 说明",
            "",
            "本报告由系统自动生成。建议开启「AI 总结」以获得更完整的洞察与建议。",
            "",
        ]
    )
    return "\n".join(lines)


def _build_appendix_md(stats: dict[str, Any]) -> str:
    """精简数据附录（Markdown）。"""
    summary = stats.get("summary") or {}
    lines = [
        "## 数据附录",
        "",
        f"- 工作区：{stats.get('workspace_id') or '—'}",
        f"- 统计范围：{_range_label(stats)}",
        f"- 范围内 / 全量：{summary.get('total_in_range')} / {summary.get('total_all')}",
        f"- 未关闭 / 已关闭：{summary.get('open')} / {summary.get('closed')}",
        "",
    ]
    return "\n".join(lines)


def compose_export_markdown(
    stats: dict[str, Any],
    ai_analysis: str | None = None,
    *,
    include_appendix: bool = True,
) -> str:
    """组装最终导出 Markdown：AI/模板正文 + 可选精简附录。"""
    ws = stats.get("workspace_id") or "—"
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    body = (ai_analysis or "").strip() or _build_fallback_narrative(stats)

    parts = [
        f"# TAPD 缺陷统计报告",
        "",
        f"> 工作区 {ws} · {_range_label(stats)} · 生成于 {now}",
        "",
        body,
    ]
    if include_appendix:
        parts.extend(["", _build_appendix_md(stats)])
    return "\n".join(parts)


def _simple_md_to_html(md: str) -> str:
    """markdown 包不可用时的简易转换（保证 PDF 仍可导出）。"""
    import html as html_lib

    lines: list[str] = []
    in_ul = False
    for raw in md.replace("\r\n", "\n").split("\n"):
        line = raw.rstrip()
        if not line:
            if in_ul:
                lines.append("</ul>")
                in_ul = False
            lines.append("<br/>")
            continue
        if line.startswith("### "):
            if in_ul:
                lines.append("</ul>")
                in_ul = False
            lines.append(f"<h3>{html_lib.escape(line[4:])}</h3>")
        elif line.startswith("## "):
            if in_ul:
                lines.append("</ul>")
                in_ul = False
            lines.append(f"<h2>{html_lib.escape(line[3:])}</h2>")
        elif line.startswith("# "):
            if in_ul:
                lines.append("</ul>")
                in_ul = False
            lines.append(f"<h1>{html_lib.escape(line[2:])}</h1>")
        elif line.startswith("> "):
            if in_ul:
                lines.append("</ul>")
                in_ul = False
            lines.append(f"<p class='meta'>{html_lib.escape(line[2:])}</p>")
        elif line.startswith("- ") or line.startswith("* "):
            if not in_ul:
                lines.append("<ul>")
                in_ul = True
            text = html_lib.escape(line[2:])
            text = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", text)
            lines.append(f"<li>{text}</li>")
        else:
            if in_ul:
                lines.append("</ul>")
                in_ul = False
            text = html_lib.escape(line)
            text = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", text)
            lines.append(f"<p>{text}</p>")
    if in_ul:
        lines.append("</ul>")
    body = "\n".join(lines)
    return f"<html><head><meta charset='utf-8'></head><body>{body}</body></html>"


def _markdown_to_html(md: str) -> str:
    try:
        import markdown as md_lib

        body = md_lib.markdown(
            md,
            extensions=["tables", "nl2br", "sane_lists"],
        )
        return f"<html><head><meta charset='utf-8'></head><body>{body}</body></html>"
    except ImportError:
        logger.warning("[tapd_export] markdown not installed, using simple HTML fallback")
        return _simple_md_to_html(md)


def _split_html_for_pages(full_html: str, max_chars: int = 12000) -> list[str]:
    """按块拆分 HTML，避免单页 insert_htmlbox 溢出。"""
    if len(full_html) <= max_chars:
        return [full_html]

    chunks: list[str] = []
    current = full_html
    markers = ["<h2", "<h1", "<hr"]
    while len(current) > max_chars:
        split_at = -1
        for marker in markers:
            idx = current.rfind(marker, 0, max_chars)
            if idx > split_at:
                split_at = idx
        if split_at <= 0:
            split_at = max_chars
        chunks.append(current[:split_at])
        current = current[split_at:]
    if current.strip():
        chunks.append(current)
    return chunks


def build_pdf_bytes(markdown: str) -> bytes:
    """Markdown → HTML → 渲染 PDF（非源码直出）。"""
    import fitz

    full_html = _markdown_to_html(markdown)
    chunks = _split_html_for_pages(full_html)
    doc = fitz.open()
    rect = fitz.Rect(45, 45, 550, 797)

    for chunk in chunks:
        page = doc.new_page(width=595, height=842)
        page.insert_htmlbox(rect, chunk, css=REPORT_CSS)

    out = doc.tobytes()
    doc.close()
    return out


def export_filename_base(stats: dict[str, Any]) -> str:
    ws = str(stats.get("workspace_id") or "tapd").strip()
    tag = datetime.now().strftime("%Y%m%d")
    days = (stats.get("range") or {}).get("days")
    range_tag = "all" if days in (None, 0) else f"{days}d"
    safe = re.sub(r"[^\w\-]", "", ws)[:20] or "tapd"
    return f"tapd_bug_stats_{safe}_{range_tag}_{tag}"


def build_excel_bytes(stats: dict[str, Any], ai_analysis: str | None = None) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill

    wb = Workbook()
    summary = stats.get("summary") or {}

    if ai_analysis and ai_analysis.strip():
        ws_ai = wb.active
        ws_ai.title = "AI分析报告"
        ws_ai["A1"] = "TAPD 缺陷统计 — AI 整理报告"
        ws_ai["A1"].font = Font(bold=True, size=14)
        ws_ai["A2"] = f"工作区 {stats.get('workspace_id')} · {_range_label(stats)}"
        ws_ai["A4"] = ai_analysis.strip()
        ws_ai["A4"].alignment = Alignment(wrap_text=True, vertical="top")
        ws_ai.column_dimensions["A"].width = 110
        ws0 = wb.create_sheet("概览")
    else:
        ws0 = wb.active
        ws0.title = "概览"

    ws0["A1"] = "TAPD 缺陷统计报告"
    ws0["A1"].font = Font(bold=True, size=14)
    rows = [
        ("工作区", stats.get("workspace_id")),
        ("统计范围", _range_label(stats)),
        ("生成时间", datetime.now().strftime("%Y-%m-%d %H:%M")),
        ("范围内缺陷", summary.get("total_in_range")),
        ("工作区总计", summary.get("total_all")),
        ("未关闭", summary.get("open")),
        ("已关闭", summary.get("closed")),
        ("未关闭占比", f"{summary.get('open_rate')}%"),
    ]
    for i, (k, v) in enumerate(rows, 3):
        ws0.cell(row=i, column=1, value=k)
        ws0.cell(row=i, column=2, value=v)

    def _dist_sheet(title: str, key: str) -> None:
        ws = wb.create_sheet(title[:31])
        ws.append(["项目", "数量", "占比(%)"])
        for cell in ws[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill("solid", fgColor="E8EAF6")
        for row in stats.get(key) or []:
            ws.append([row.get("label"), row.get("count"), row.get("percent")])

    _dist_sheet("按状态", "by_status")
    _dist_sheet("按优先级", "by_priority")
    _dist_sheet("按处理人", "by_owner")
    _dist_sheet("按报告人", "by_reporter")

    ws_bug = wb.create_sheet("缺陷明细")
    ws_bug.append(["ID", "标题", "状态", "优先级", "处理人", "报告人", "创建时间", "链接"])
    for cell in ws_bug[1]:
        cell.font = Font(bold=True)
    for b in stats.get("recent_bugs") or []:
        ws_bug.append(
            [
                b.get("id"),
                b.get("title"),
                b.get("status_label") or b.get("status"),
                b.get("priority_label"),
                b.get("current_owner"),
                b.get("reporter"),
                b.get("created"),
                b.get("url"),
            ]
        )

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _md_to_plain_slides(md: str) -> list[tuple[str, list[str]]]:
    """从 Markdown 提取幻灯片标题与要点。"""
    slides: list[tuple[str, list[str]]] = []
    current_title = "分析要点"
    current_lines: list[str] = []

    for line in md.splitlines():
        s = line.strip()
        if not s:
            continue
        if s.startswith("#"):
            if current_lines:
                slides.append((current_title, current_lines))
                current_lines = []
            current_title = re.sub(r"^#+\s*", "", s).strip() or "分析要点"
            continue
        if s.startswith(("- ", "* ", "• ")):
            current_lines.append(re.sub(r"^[-*•]\s*", "", s))
        elif re.match(r"^\d+\.", s):
            current_lines.append(s)
        elif len(s) > 4:
            current_lines.append(s)

    if current_lines:
        slides.append((current_title, current_lines))
    return slides


def build_ppt_bytes(stats: dict[str, Any], ai_analysis: str | None = None) -> bytes:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    summary = stats.get("summary") or {}

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TAPD 缺陷统计报告"
    slide.placeholders[1].text = (
        f"工作区 {stats.get('workspace_id')}\n"
        f"{_range_label(stats)}\n"
        f"{datetime.now().strftime('%Y-%m-%d %H:%M')}"
    )

    if ai_analysis and ai_analysis.strip():
        for title, bullets in _md_to_plain_slides(ai_analysis)[:8]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title[:80]
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, text in enumerate(bullets[:6]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = text[:300]
                p.font.size = Pt(16)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "核心指标"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for line in [
            f"范围内缺陷：{summary.get('total_in_range', '—')}",
            f"未关闭：{summary.get('open', '—')}（{summary.get('open_rate', '—')}%）",
            f"已关闭：{summary.get('closed', '—')}",
        ]:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "按状态分布"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for row in (stats.get("by_status") or [])[:8]:
        p = tf.add_paragraph()
        p.text = f"{row.get('label')}：{row.get('count')}（{row.get('percent')}%）"
        p.font.size = Pt(16)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def export_stats_report(
    stats: dict[str, Any],
    *,
    fmt: str,
    ai_analysis: str | None = None,
) -> tuple[bytes, str, str]:
    """返回 (文件字节, mime_type, filename)。"""
    ext = (fmt or "md").strip().lower()
    if ext not in EXPORT_FORMATS:
        raise ValueError(f"不支持的格式: {ext}")

    base = export_filename_base(stats)
    narrative_md = compose_export_markdown(stats, ai_analysis, include_appendix=(ext == "md"))

    if ext == "md":
        content = narrative_md.encode("utf-8")
        return content, "text/markdown; charset=utf-8", f"{base}.md"
    if ext == "excel":
        return (
            build_excel_bytes(stats, ai_analysis),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            f"{base}.xlsx",
        )
    if ext == "pdf":
        return build_pdf_bytes(narrative_md), "application/pdf", f"{base}.pdf"
    if ext == "ppt":
        return (
            build_ppt_bytes(stats, ai_analysis),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            f"{base}.pptx",
        )
    raise ValueError(f"不支持的格式: {ext}")
